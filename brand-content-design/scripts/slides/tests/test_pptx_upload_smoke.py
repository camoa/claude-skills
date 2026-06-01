"""End-to-end smoke test for the PPTX-import path (Drive convert-on-upload).

Skipped automatically when:

* Google credentials env vars are absent, OR
* ``python-pptx`` is not installed.

When run, this test:

  1. Generates a minimal one-slide PPTX with python-pptx (no external file).
  2. Authenticates via :func:`slides.auth.build_services`.
  3. Uploads the PPTX with ``mimeType:
     application/vnd.google-apps.presentation`` — Drive's OOXML importer
     converts it to a native Slides deck on upload.
  4. Reads the file back and asserts ``mimeType`` is the Slides type
     (proving the import succeeded).
  5. Trashes the deck in ``finally`` so the run leaves no orphaned
     artifacts.

Validates the canonical PPTX-import Slides rendering path the
brand-content-design pivot moved to in 2026-05-26.
"""

from __future__ import annotations

import io
import os
import sys

import pytest


CREDS_PRESENT = bool(
    os.environ.get("BCD_SLIDES_SA_KEY_FILE")
    or (
        os.environ.get("BCD_SLIDES_OAUTH_CLIENT_ID")
        and os.environ.get("BCD_SLIDES_OAUTH_CLIENT_SECRET")
        and os.environ.get("BCD_SLIDES_OAUTH_REFRESH_TOKEN")
    )
)

try:
    import pptx  # noqa: F401 — feature check only
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


pytestmark = [
    pytest.mark.skipif(not CREDS_PRESENT, reason="E2E creds not configured"),
    pytest.mark.skipif(
        not PPTX_AVAILABLE, reason="python-pptx not installed for smoke"
    ),
]


SMOKE_LITERAL = "pptx-import smoke 2026-05-26"
GOOGLE_SLIDES_MIMETYPE = "application/vnd.google-apps.presentation"
PPTX_MIMETYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)


def _redact(file_id: str) -> str:
    if len(file_id) <= 8:
        return "***"
    return f"{file_id[:4]}...{file_id[-4:]}"


def _make_minimal_pptx(tmp_path) -> str:
    """Generate a tiny PPTX on disk and return the absolute path."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title-only layout
    title = slide.shapes.title
    title.text = SMOKE_LITERAL
    out = tmp_path / "smoke.pptx"
    prs.save(str(out))
    return str(out)


def test_pptx_uploads_and_converts_to_slides(tmp_path):
    from googleapiclient.http import MediaFileUpload

    from slides.auth import build_services, resolve_mode

    mode = resolve_mode()
    print(f"[pptx-smoke] auth mode: {mode}", file=sys.stderr)

    pptx_path = _make_minimal_pptx(tmp_path)
    _slides_svc, drive_svc = build_services()

    media = MediaFileUpload(pptx_path, mimetype=PPTX_MIMETYPE, resumable=False)
    created = (
        drive_svc.files()
        .create(
            body={
                "name": f"camoa-skills pptx-smoke {SMOKE_LITERAL}",
                "mimeType": GOOGLE_SLIDES_MIMETYPE,
            },
            media_body=media,
            fields="id, mimeType",
        )
        .execute()
    )
    file_id = created["id"]
    print(f"[pptx-smoke] uploaded+converted: {_redact(file_id)}", file=sys.stderr)

    try:
        # Read back: mimeType should now be the Slides type (Drive converted).
        meta = (
            drive_svc.files()
            .get(fileId=file_id, fields="id, mimeType, name")
            .execute()
        )
        assert meta["mimeType"] == GOOGLE_SLIDES_MIMETYPE, (
            f"Expected mimeType {GOOGLE_SLIDES_MIMETYPE}, got "
            f"{meta['mimeType']} — Drive did not convert on upload"
        )
        print(
            f"[pptx-smoke] verified mimeType: {meta['mimeType']}",
            file=sys.stderr,
        )
    finally:
        drive_svc.files().update(fileId=file_id, body={"trashed": True}).execute()
        print(f"[pptx-smoke] trashed: {_redact(file_id)}", file=sys.stderr)
