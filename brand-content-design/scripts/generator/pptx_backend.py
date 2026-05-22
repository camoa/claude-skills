#!/usr/bin/env python3
"""
PPTX backend for the template generator — python-pptx.

`PptxCanvas` subclasses `DeckCanvas` and renders to a native, editable PowerPoint
deck — real shapes and text, not an image dump. It is template-agnostic: the
template's font table (a `{logical_name: FontSpec}` map) is passed to the
constructor.

Text-as-block (contract C2): the spec passes an unwrapped string + a box width;
this backend emits ONE word-wrapped text box for it. A wrapped headline is one
editable paragraph, not N stacked boxes — the structural fix for the
multi-line-text bug. Because (x, y) is already the box top-left there is no
baseline reverse-engineering and no ascent fudge.

Coordinate conversion: the spec is authored in a 1920x1080 px space, origin
bottom-left. PPTX is top-left, in EMU, on a 10 x 5.625 in slide:
  EMU_PER_PX = 914400 / 192 ;  point size = px * 0.375
  shapes (bottom-left corner):  pptx_top = HEIGHT - y - h
  text box (top-left corner)  :  pptx_top = HEIGHT - y
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn, nsdecls

from deck_canvas import DeckCanvas, WIDTH, HEIGHT

EMU_PER_PX = 914400 / 192.0   # 1920 px -> 10 in
PT_PER_PX = 0.375             # 1920 px -> 720 pt

_ALIGN = {'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT,
          'left': PP_ALIGN.LEFT}
_ANCHOR = {'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE,
           'bottom': MSO_ANCHOR.BOTTOM}


def _emu(px):
    return Emu(int(round(px * EMU_PER_PX)))


def _rgb(hexstr):
    return RGBColor.from_string(hexstr.lstrip('#').upper())


class PptxCanvas(DeckCanvas):
    """Renders the abstract DeckCanvas primitives to native PowerPoint shapes."""

    def __init__(self, out, fonts):
        """fonts: {logical_name: FontSpec} — the PPTX family name + bold flag
        per logical font are read from each FontSpec."""
        self.out = out
        self.fonts = fonts
        self.prs = Presentation()
        self.prs.slide_width = _emu(WIDTH)
        self.prs.slide_height = _emu(HEIGHT)
        self._blank = self.prs.slide_layouts[6]
        self.slide = None

    def start_slide(self):
        self.slide = self.prs.slides.add_slide(self._blank)

    def save(self):
        self.prs.save(self.out)

    # spec (x, y) is a shape's bottom-left corner, y up -> PPTX top-left EMU box
    def _box(self, x, y, w, h):
        return _emu(x), _emu(HEIGHT - y - h), _emu(w), _emu(h)

    @staticmethod
    def _strip_shadow(shape):
        try:
            shape.shadow.inherit = False
        except Exception:
            pass

    def _solid(self, color):
        self._rect(0, 0, WIDTH, HEIGHT, color)

    def _gradient(self, stops):
        left, top, w, h = self._box(0, 0, WIDTH, HEIGHT)
        shp = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
        self._strip_shadow(shp)
        shp.line.fill.background()
        shp.fill.solid()  # places <a:solidFill> in the schema-correct slot
        spPr = shp._element.spPr
        solid = spPr.find(qn('a:solidFill'))
        gs = ''.join(
            f'<a:gs pos="{int(round(p * 100000))}">'
            f'<a:srgbClr val="{c.lstrip("#").upper()}"/></a:gs>'
            for p, c in stops)
        # diagonal top-left to bottom-right (45 deg)
        grad = parse_xml(
            '<a:gradFill %s rotWithShape="1">'
            '<a:gsLst>%s</a:gsLst>'
            '<a:lin ang="2700000" scaled="1"/>'
            '</a:gradFill>' % (nsdecls('a'), gs))
        solid.addprevious(grad)
        spPr.remove(solid)

    def _rect(self, x, y, w, h, color):
        left, top, w, h = self._box(x, y, w, h)
        shp = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
        self._strip_shadow(shp)
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(color)
        shp.line.fill.background()

    def _round_rect(self, x, y, w, h, r, fill, stroke, stroke_w):
        left, top, ww, hh = self._box(x, y, w, h)
        shp = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, ww, hh)
        self._strip_shadow(shp)
        try:
            shp.adjustments[0] = max(0.0, min(0.5, r / float(min(w, h))))
        except Exception:
            pass
        if fill:
            shp.fill.solid()
            shp.fill.fore_color.rgb = _rgb(fill)
        else:
            shp.fill.background()
        if stroke:
            shp.line.color.rgb = _rgb(stroke)
            shp.line.width = _emu(stroke_w)
        else:
            shp.line.fill.background()

    def _circle(self, cx, cy, r, fill, stroke, stroke_w):
        left, top, ww, hh = self._box(cx - r, cy - r, 2 * r, 2 * r)
        shp = self.slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, ww, hh)
        self._strip_shadow(shp)
        if fill:
            shp.fill.solid()
            shp.fill.fore_color.rgb = _rgb(fill)
        else:
            shp.fill.background()
        if stroke:
            shp.line.color.rgb = _rgb(stroke)
            shp.line.width = _emu(stroke_w)
        else:
            shp.line.fill.background()

    def _line(self, x1, y1, x2, y2, color, w):
        cn = self.slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            _emu(x1), _emu(HEIGHT - y1), _emu(x2), _emu(HEIGHT - y2))
        self._strip_shadow(cn)
        cn.line.color.rgb = _rgb(color)
        cn.line.width = _emu(w)

    def _text_block(self, s, x, y, w, h, font, size, color, align, valign):
        """One word-wrapped text box. (x, y) is the box top-left in spec space;
        the box top in PPTX is HEIGHT - y. A multi-line string wraps inside this
        single box — it is never split into per-line boxes."""
        if not s:
            return
        spec = self.fonts[font]
        tb = self.slide.shapes.add_textbox(
            _emu(x), _emu(HEIGHT - y), _emu(w), _emu(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = _ANCHOR.get(valign, MSO_ANCHOR.TOP)
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = _ALIGN.get(align, PP_ALIGN.LEFT)
        try:
            p.line_spacing = 1.25
        except Exception:
            pass
        run = p.add_run()
        run.text = s
        f = run.font
        f.name = spec.pptx
        f.bold = bool(spec.bold)
        f.size = Pt(size * PT_PER_PX)
        f.color.rgb = _rgb(color)

    def _image(self, path, x, y, w, h):
        left, top, ww, hh = self._box(x, y, w, h)
        try:
            self.slide.shapes.add_picture(path, left, top, width=ww, height=hh)
        except Exception:
            pass
